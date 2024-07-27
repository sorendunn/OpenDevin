"""
agentskills.py

This module provides various file manipulation skills for the OpenDevin agent.

Functions:
- open_file(path: str, line_number: int | None = 1, context_lines: int = 100): Opens a file and optionally moves to a specific line.
- goto_line(line_number): Moves the window to show the specified line number.
- scroll_down(): Moves the window down by the number of lines specified in WINDOW.
- scroll_up(): Moves the window up by the number of lines specified in WINDOW.
- create_file(filename): Creates and opens a new file with the given name.
- search_dir(search_term, dir_path='./'): Searches for a term in all files in the specified directory.
- search_file(search_term, file_path=None): Searches for a term in the specified file or the currently open file.
- find_file(file_name, dir_path='./'): Finds all files with the given name in the specified directory.
- edit_file(file_name: str, to_replace: str, new_content: str): Replaces lines in a file with the given content.
- insert_content_at_line(file_name: str, line_number: int, content: str): Inserts given content at the specified line number in a file.
"""

import base64
import functools
import os
import re
import shutil
import subprocess
import tempfile
from inspect import signature
from typing import Optional

import docx
import PyPDF2
from openai import OpenAI
from pptx import Presentation
from pylatexenc.latex2text import LatexNodes2Text

CURRENT_FILE: str | None = None
CURRENT_LINE = 1
WINDOW = 100

ENABLE_AUTO_LINT = os.getenv('ENABLE_AUTO_LINT', 'false').lower() == 'true'

# This is also used in unit tests!
MSG_FILE_UPDATED = '[File updated. Please review the changes and make sure they are correct (correct indentation, no duplicate lines, etc). Edit the file again if necessary.]'

# OPENAI
OPENAI_API_KEY = os.getenv(
    'OPENAI_API_KEY', os.getenv('SANDBOX_ENV_OPENAI_API_KEY', '')
)
OPENAI_BASE_URL = os.getenv('OPENAI_BASE_URL', 'https://api.openai.com/v1')
OPENAI_MODEL = os.getenv('OPENAI_MODEL', 'gpt-4o-mini-2024-07-18')
MAX_TOKEN = os.getenv('MAX_TOKEN', 300)

OPENAI_PROXY = f'{OPENAI_BASE_URL}/chat/completions'

client = OpenAI(api_key=OPENAI_API_KEY, base_url=OPENAI_BASE_URL)


# Define the decorator using the functionality of UpdatePwd
def update_pwd_decorator(func):
    @functools.wraps(func)
    def wrapper(*args, **kwargs):
        old_pwd = os.getcwd()
        jupyter_pwd = os.environ.get('JUPYTER_PWD', None)
        if jupyter_pwd:
            os.chdir(jupyter_pwd)
        try:
            return func(*args, **kwargs)
        finally:
            os.chdir(old_pwd)

    return wrapper


def _is_valid_filename(file_name) -> bool:
    if not file_name or not isinstance(file_name, str) or not file_name.strip():
        return False
    invalid_chars = '<>:"/\\|?*'
    if os.name == 'nt':  # Windows
        invalid_chars = '<>:"/\\|?*'
    elif os.name == 'posix':  # Unix-like systems
        invalid_chars = '\0'

    for char in invalid_chars:
        if char in file_name:
            return False
    return True


def _is_valid_path(path) -> bool:
    if not path or not isinstance(path, str):
        return False
    try:
        return os.path.exists(os.path.normpath(path))
    except PermissionError:
        return False


def _create_paths(file_name) -> bool:
    try:
        dirname = os.path.dirname(file_name)
        if dirname:
            os.makedirs(dirname, exist_ok=True)
        return True
    except PermissionError:
        return False


def _check_current_file(file_path: str | None = None) -> bool:
    global CURRENT_FILE
    if not file_path:
        file_path = CURRENT_FILE
    if not file_path or not os.path.isfile(file_path):
        raise ValueError('No file open. Use the open_file function first.')
    return True


def _clamp(value, min_value, max_value):
    return max(min_value, min(value, max_value))


def _lint_file(file_path: str) -> tuple[Optional[str], Optional[int]]:
    """
    Lint the file at the given path and return a tuple with a boolean indicating if there are errors,
    and the line number of the first error, if any.

    Returns:
        tuple[str, Optional[int]]: (lint_error, first_error_line_number)
    """

    if file_path.endswith('.py'):
        # Define the flake8 command with selected error codes
        def _command_fn(executable):
            return [
                executable,
                '--isolated',
                '--select=F821,F822,F831,E112,E113,E999,E902',
                file_path,
            ]

        if os.path.exists('/opendevin/miniforge3/bin/flake8'):
            # when this function is called from the docker sandbox,
            # the flake8 command is available at /opendevin/miniforge3/bin/flake8
            executable = '/opendevin/miniforge3/bin/flake8'
        else:
            executable = 'flake8'

        command = _command_fn(executable)
        result = subprocess.run(
            command,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
        )
        if result.returncode == 0:
            # Linting successful. No issues found.
            return None, None

        # Extract the line number from the first error message
        error_message = result.stdout.decode().strip()
        lint_error = 'ERRORS:\n' + error_message
        first_error_line = None
        for line in error_message.splitlines(True):
            if line.strip():
                # The format of the error message is: <filename>:<line>:<column>: <error code> <error message>
                parts = line.split(':')
                if len(parts) >= 2:
                    try:
                        first_error_line = int(parts[1])
                        break
                    except ValueError:
                        # Not a valid line number, continue to the next line
                        continue

        return lint_error, first_error_line

    # Not a python file, skip linting
    return None, None


def _print_window(file_path, targeted_line, window, return_str=False):
    global CURRENT_LINE
    _check_current_file(file_path)
    with open(file_path) as file:
        content = file.read()

        # Ensure the content ends with a newline character
        if not content.endswith('\n'):
            content += '\n'

        lines = content.splitlines(True)  # Keep all line ending characters
        total_lines = len(lines)

        # cover edge cases
        CURRENT_LINE = _clamp(targeted_line, 1, total_lines)
        half_window = max(1, window // 2)

        # Ensure at least one line above and below the targeted line
        start = max(1, CURRENT_LINE - half_window)
        end = min(total_lines, CURRENT_LINE + half_window)

        # Adjust start and end to ensure at least one line above and below
        if start == 1:
            end = min(total_lines, start + window - 1)
        if end == total_lines:
            start = max(1, end - window + 1)

        output = ''

        # only display this when there's at least one line above
        if start > 1:
            output += f'({start - 1} more lines above)\n'
        for i in range(start, end + 1):
            _new_line = f'{i}|{lines[i-1]}'
            if not _new_line.endswith('\n'):
                _new_line += '\n'
            output += _new_line
        if end < total_lines:
            output += f'({total_lines - end} more lines below)\n'
        output = output.rstrip()

        if return_str:
            return output
        else:
            print(output)


def _cur_file_header(current_file, total_lines) -> str:
    if not current_file:
        return ''
    return f'[File: {os.path.abspath(current_file)} ({total_lines} lines total)]\n'


@update_pwd_decorator
def open_file(
    path: str, line_number: int | None = 1, context_lines: int | None = 100
) -> None:
    """
    Opens the file at the given path in the editor. If line_number is provided, the window will be moved to include that line.
    It only shows the first 100 lines by default! Max `context_lines` supported is 2000, use `scroll up/down`
    to view the file if you want to see more.

    Args:
        path: str: The path to the file to open, preferred absolute path.
        line_number: int | None = 1: The line number to move to. Defaults to 1.
        context_lines: int | None = 100: Only shows this number of lines in the context window (usually from line 1), with line_number as the center (if possible). Defaults to 100.
    """
    global CURRENT_FILE, CURRENT_LINE, WINDOW

    if not os.path.isfile(path):
        raise FileNotFoundError(f'File {path} not found')

    CURRENT_FILE = os.path.abspath(path)
    with open(CURRENT_FILE) as file:
        total_lines = max(1, sum(1 for _ in file))

    if not isinstance(line_number, int) or line_number < 1 or line_number > total_lines:
        raise ValueError(f'Line number must be between 1 and {total_lines}')
    CURRENT_LINE = line_number

    # Override WINDOW with context_lines
    if context_lines is None or context_lines < 1:
        context_lines = 100
    WINDOW = _clamp(context_lines, 1, 2000)

    output = _cur_file_header(CURRENT_FILE, total_lines)
    output += _print_window(CURRENT_FILE, CURRENT_LINE, WINDOW, return_str=True)
    print(output)


@update_pwd_decorator
def goto_line(line_number: int) -> None:
    """
    Moves the window to show the specified line number.

    Args:
        line_number: int: The line number to move to.
    """
    global CURRENT_FILE, CURRENT_LINE, WINDOW
    _check_current_file()

    with open(str(CURRENT_FILE)) as file:
        total_lines = max(1, sum(1 for _ in file))
    if not isinstance(line_number, int) or line_number < 1 or line_number > total_lines:
        raise ValueError(f'Line number must be between 1 and {total_lines}')

    CURRENT_LINE = _clamp(line_number, 1, total_lines)

    output = _cur_file_header(CURRENT_FILE, total_lines)
    output += _print_window(CURRENT_FILE, CURRENT_LINE, WINDOW, return_str=True)
    print(output)


@update_pwd_decorator
def scroll_down() -> None:
    """Moves the window down by 100 lines.

    Args:
        None
    """
    global CURRENT_FILE, CURRENT_LINE, WINDOW
    _check_current_file()

    with open(str(CURRENT_FILE)) as file:
        total_lines = max(1, sum(1 for _ in file))
    CURRENT_LINE = _clamp(CURRENT_LINE + WINDOW, 1, total_lines)
    output = _cur_file_header(CURRENT_FILE, total_lines)
    output += _print_window(CURRENT_FILE, CURRENT_LINE, WINDOW, return_str=True)
    print(output)


@update_pwd_decorator
def scroll_up() -> None:
    """Moves the window up by 100 lines.

    Args:
        None
    """
    global CURRENT_FILE, CURRENT_LINE, WINDOW
    _check_current_file()

    with open(str(CURRENT_FILE)) as file:
        total_lines = max(1, sum(1 for _ in file))
    CURRENT_LINE = _clamp(CURRENT_LINE - WINDOW, 1, total_lines)
    output = _cur_file_header(CURRENT_FILE, total_lines)
    output += _print_window(CURRENT_FILE, CURRENT_LINE, WINDOW, return_str=True)
    print(output)


@update_pwd_decorator
def create_file(filename: str) -> None:
    """Creates and opens a new file with the given name.

    Args:
        filename: str: The name of the file to create.
    """
    if os.path.exists(filename):
        raise FileExistsError(f"File '{filename}' already exists.")

    with open(filename, 'w') as file:
        file.write('\n')

    open_file(filename)
    print(f'[File {filename} created.]')


LINTER_ERROR_MSG = '[Your proposed edit has introduced new syntax error(s). Please understand the errors and retry your edit command.]\n'


def _edit_or_insert_file(
    file_name: str,
    start: int | None = None,
    end: int | None = None,
    content: str = '',
    is_insert: bool = False,
) -> str:
    """Internal method to handle common logic for edit_/append_file methods.

    Args:
        file_name: str: The name of the file to edit or append to.
        start: int | None = None: The start line number for editing. Ignored if is_append is True.
        end: int | None = None: The end line number for editing. Ignored if is_append is True.
        content: str: The content to replace the lines with or to append.
        is_insert: bool = False: Whether to insert content at the given line number instead of editing.
    """
    ret_str = ''
    global CURRENT_FILE, CURRENT_LINE, WINDOW

    ERROR_MSG = f'[Error editing file {file_name}. Please confirm the file is correct.]'
    ERROR_MSG_SUFFIX = (
        'Your changes have NOT been applied. Please fix your edit command and try again.\n'
        'You either need to 1) Open the correct file and try again or 2) Specify the correct line number arguments.\n'
        'DO NOT re-run the same failed edit command. Running it again will lead to the same error.'
    )

    if not _is_valid_filename(file_name):
        raise FileNotFoundError('Invalid file name.')

    if not _is_valid_path(file_name):
        raise FileNotFoundError('Invalid path or file name.')

    if not _create_paths(file_name):
        raise PermissionError('Could not access or create directories.')

    if not os.path.isfile(file_name):
        raise FileNotFoundError(f'File {file_name} not found.')

    # Use a temporary file to write changes
    content = str(content or '')
    temp_file_path = ''
    src_abs_path = os.path.abspath(file_name)
    first_error_line = None
    try:
        # Create a temporary file
        with tempfile.NamedTemporaryFile('w', delete=False) as temp_file:
            temp_file_path = temp_file.name

            # Read the original file and check if empty and for a trailing newline
            with open(file_name) as original_file:
                lines = original_file.readlines()

            if is_insert:
                if len(lines) == 0:
                    new_lines = [
                        content + '\n' if not content.endswith('\n') else content
                    ]
                elif start is not None:
                    if len(lines) == 1 and lines[0].strip() == '':
                        # if the file is empty with only 1 line
                        lines = ['\n']
                    new_lines = (
                        lines[: start - 1]
                        + [content + '\n' if not content.endswith('\n') else content]
                        + lines[start - 1 :]
                    )
                else:
                    assert start is None
                    ret_str += (
                        f'{ERROR_MSG}\n'
                        f'Invalid line number: {start}. Line numbers must be between 1 and {len(lines)} (inclusive).\n'
                        f'{ERROR_MSG_SUFFIX}'
                    ) + '\n'

                content = ''.join(new_lines)
            else:
                # Handle cases where start or end are None
                if start is None:
                    start = 1  # Default to the beginning
                if end is None:
                    end = len(lines)  # Default to the end
                # Check arguments
                if not (1 <= start <= len(lines)):
                    ret_str += (
                        f'{ERROR_MSG}\n'
                        f'Invalid start line number: {start}. Line numbers must be between 1 and {len(lines)} (inclusive).\n'
                        f'{ERROR_MSG_SUFFIX}'
                    ) + '\n'
                if not (1 <= end <= len(lines)):
                    ret_str += (
                        f'{ERROR_MSG}\n'
                        f'Invalid end line number: {end}. Line numbers must be between 1 and {len(lines)} (inclusive).\n'
                        f'{ERROR_MSG_SUFFIX}'
                    ) + '\n'
                if start > end:
                    ret_str += (
                        f'{ERROR_MSG}\n'
                        f'Invalid line range: {start}-{end}. Start must be less than or equal to end.\n'
                        f'{ERROR_MSG_SUFFIX}'
                    ) + '\n'
                if not content.endswith('\n'):
                    content += '\n'
                content_lines = content.splitlines(True)
                new_lines = lines[: start - 1] + content_lines + lines[end:]
                content = ''.join(new_lines)

            if not content.endswith('\n'):
                content += '\n'

            # Write the new content to the temporary file
            temp_file.write(content)

        # Replace the original file with the temporary file atomically
        shutil.move(temp_file_path, src_abs_path)

        # Handle linting
        if ENABLE_AUTO_LINT:
            # BACKUP the original file
            original_file_backup_path = os.path.join(
                os.path.dirname(file_name),
                f'.backup.{os.path.basename(file_name)}',
            )
            with open(original_file_backup_path, 'w') as f:
                f.writelines(lines)

            lint_error, first_error_line = _lint_file(file_name)
            if lint_error is not None:
                if first_error_line is not None:
                    CURRENT_LINE = int(first_error_line)
                ret_str += LINTER_ERROR_MSG
                ret_str += lint_error + '\n'

                ret_str += '[This is how your edit would have looked if applied]\n'
                ret_str += '-------------------------------------------------\n'
                ret_str += (
                    _print_window(file_name, CURRENT_LINE, 10, return_str=True) + '\n'
                )
                ret_str += '-------------------------------------------------\n\n'

                ret_str += '[This is the original code before your edit]\n'
                ret_str += '-------------------------------------------------\n'
                ret_str += (
                    _print_window(
                        original_file_backup_path, CURRENT_LINE, 10, return_str=True
                    )
                    + '\n'
                )
                ret_str += '-------------------------------------------------\n'

                ret_str += (
                    'Your changes have NOT been applied. Please fix your edit command and try again.\n'
                    'You either need to 1) Specify the correct start/end line arguments or 2) Correct your edit code.\n'
                    'DO NOT re-run the same failed edit command. Running it again will lead to the same error.'
                )

                # recover the original file
                with open(original_file_backup_path) as fin, open(
                    file_name, 'w'
                ) as fout:
                    fout.write(fin.read())
                os.remove(original_file_backup_path)
                return ret_str

    except FileNotFoundError as e:
        ret_str += f'File not found: {e}\n'
    except IOError as e:
        ret_str += f'An error occurred while handling the file: {e}\n'
    except ValueError as e:
        ret_str += f'Invalid input: {e}\n'
    except Exception as e:
        # Clean up the temporary file if an error occurs
        if temp_file_path and os.path.exists(temp_file_path):
            os.remove(temp_file_path)
        print(f'An unexpected error occurred: {e}')
        raise e

    # Update the file information and print the updated content
    with open(file_name, 'r', encoding='utf-8') as file:
        n_total_lines = max(1, len(file.readlines()))
    if first_error_line is not None and int(first_error_line) > 0:
        CURRENT_LINE = first_error_line
    else:
        CURRENT_LINE = start or n_total_lines or 1
    ret_str += f'[File: {os.path.abspath(file_name)} ({n_total_lines} lines total after edit)]\n'
    CURRENT_FILE = file_name
    ret_str += _print_window(CURRENT_FILE, CURRENT_LINE, WINDOW, return_str=True) + '\n'
    ret_str += MSG_FILE_UPDATED
    return ret_str


@update_pwd_decorator
def edit_file(file_name: str, to_replace: str, new_content: str) -> None:
    """Edit a file. This will search for `to_replace` in the given file and replace it with `new_content`.

    Every *to_replace* must *EXACTLY MATCH* the existing source code, character for character, including all comments, docstrings, etc.

    Include enough lines to make code in `to_replace` unique. `to_replace` should NOT be empty.
    `edit_file` will only replace the *first* matching occurrences.

    For example, given a file "/workspace/example.txt" with the following content:
    ```
    line 1
    line 2
    line 2
    line 3
    ```

    EDITING: If you want to replace the second occurrence of "line 2", you can make `to_replace` unique:

    edit_file(
        '/workspace/example.txt',
        to_replace='line 2\nline 3',
        new_content='new line\nline 3',
    )

    This will replace only the second "line 2" with "new line". The first "line 2" will remain unchanged.

    The resulting file will be:
    ```
    line 1
    line 2
    new line
    line 3
    ```

    REMOVAL: If you want to remove "line 2" and "line 3", you can set `new_content` to an empty string:

    edit_file(
        '/workspace/example.txt',
        to_replace='line 2\nline 3',
        new_content='',
    )

    Args:
        file_name: str: The name of the file to edit.
        to_replace: str: The content to search for and replace.
        new_content: str: The new content to replace the old content with.
    """
    # FIXME: support replacing *all* occurrences
    if to_replace.strip() == '':
        raise ValueError('`to_replace` must not be empty.')

    # search for `to_replace` in the file
    # if found, replace it with `new_content`
    # if not found, perform a fuzzy search to find the closest match and replace it with `new_content`
    with open(file_name, 'r') as file:
        file_content = file.read()

    start = file_content.find(to_replace)
    if start != -1:
        # Convert start from index to line number
        start_line_number = file_content[:start].count('\n') + 1
        end_line_number = start_line_number + len(to_replace.splitlines()) - 1
    else:

        def _fuzzy_transform(s: str) -> str:
            # remove all space except newline
            return re.sub(r'[^\S\n]+', '', s)

        # perform a fuzzy search (remove all spaces except newlines)
        to_replace_fuzzy = _fuzzy_transform(to_replace)
        file_content_fuzzy = _fuzzy_transform(file_content)
        # find the closest match
        start = file_content_fuzzy.find(to_replace_fuzzy)
        if start == -1:
            print(
                f'[No exact match found in {file_name} for\n```\n{to_replace}\n```\n]'
            )
            return
        # Convert start from index to line number for fuzzy match
        start_line_number = file_content_fuzzy[:start].count('\n') + 1
        end_line_number = start_line_number + len(to_replace.splitlines()) - 1

    ret_str = _edit_or_insert_file(
        file_name,
        start=start_line_number,
        end=end_line_number,
        content=new_content,
        is_insert=False,
    )
    # lint_error = bool(LINTER_ERROR_MSG in ret_str)
    # TODO: automatically tries to fix linter error (maybe involve some static analysis tools on the location near the edit to figure out indentation)
    print(ret_str)


@update_pwd_decorator
def insert_content_at_line(file_name: str, line_number: int, content: str) -> None:
    """Insert content at the given line number in a file.
    This will NOT modify the content of the lines before OR after the given line number.

    For example, if the file has the following content:
    ```
    line 1
    line 2
    line 3
    ```
    and you call `insert_content_at_line('file.txt', 2, 'new line')`, the file will be updated to:
    ```
    line 1
    new line
    line 2
    line 3
    ```

    Args:
        file_name: str: The name of the file to edit.
        line_number: int: The line number (starting from 1) to insert the content after.
        content: str: The content to insert.
    """
    ret_str = _edit_or_insert_file(
        file_name, start=line_number, end=line_number, content=content, is_insert=True
    )
    print(ret_str)


@update_pwd_decorator
def search_dir(search_term: str, dir_path: str = './') -> None:
    """Searches for search_term in all files in dir. If dir is not provided, searches in the current directory.

    Args:
        search_term: str: The term to search for.
        dir_path: Optional[str]: The path to the directory to search.
    """
    if not os.path.isdir(dir_path):
        raise FileNotFoundError(f'Directory {dir_path} not found')
    matches = []
    for root, _, files in os.walk(dir_path):
        for file in files:
            if file.startswith('.'):
                continue
            file_path = os.path.join(root, file)
            with open(file_path, 'r', errors='ignore') as f:
                for line_num, line in enumerate(f, 1):
                    if search_term in line:
                        matches.append((file_path, line_num, line.strip()))

    if not matches:
        print(f'No matches found for "{search_term}" in {dir_path}')
        return

    num_matches = len(matches)
    num_files = len(set(match[0] for match in matches))

    if num_files > 100:
        print(
            f'More than {num_files} files matched for "{search_term}" in {dir_path}. Please narrow your search.'
        )
        return

    print(f'[Found {num_matches} matches for "{search_term}" in {dir_path}]')
    for file_path, line_num, line in matches:
        print(f'{file_path} (Line {line_num}): {line}')
    print(f'[End of matches for "{search_term}" in {dir_path}]')


@update_pwd_decorator
def search_file(search_term: str, file_path: Optional[str] = None) -> None:
    """Searches for search_term in file. If file is not provided, searches in the current open file.

    Args:
        search_term: str: The term to search for.
        file_path: Optional[str]: The path to the file to search.
    """
    global CURRENT_FILE
    if file_path is None:
        file_path = CURRENT_FILE
    if file_path is None:
        raise FileNotFoundError(
            'No file specified or open. Use the open_file function first.'
        )
    if not os.path.isfile(file_path):
        raise FileNotFoundError(f'File {file_path} not found')

    matches = []
    with open(file_path) as file:
        for i, line in enumerate(file, 1):
            if search_term in line:
                matches.append((i, line.strip()))

    if matches:
        print(f'[Found {len(matches)} matches for "{search_term}" in {file_path}]')
        for match in matches:
            print(f'Line {match[0]}: {match[1]}')
        print(f'[End of matches for "{search_term}" in {file_path}]')
    else:
        print(f'[No matches found for "{search_term}" in {file_path}]')


@update_pwd_decorator
def find_file(file_name: str, dir_path: str = './') -> None:
    """Finds all files with the given name in the specified directory.

    Args:
        file_name: str: The name of the file to find.
        dir_path: Optional[str]: The path to the directory to search.
    """
    if not os.path.isdir(dir_path):
        raise FileNotFoundError(f'Directory {dir_path} not found')

    matches = []
    for root, _, files in os.walk(dir_path):
        for file in files:
            if file_name in file:
                matches.append(os.path.join(root, file))

    if matches:
        print(f'[Found {len(matches)} matches for "{file_name}" in {dir_path}]')
        for match in matches:
            print(f'{match}')
        print(f'[End of matches for "{file_name}" in {dir_path}]')
    else:
        print(f'[No matches found for "{file_name}" in {dir_path}]')


@update_pwd_decorator
def parse_pdf(file_path: str) -> None:
    """Parses the content of a PDF file and prints it.

    Args:
        file_path: str: The path to the file to open.
    """
    print(f'[Reading PDF file from {file_path}]')
    content = PyPDF2.PdfReader(file_path)
    text = ''
    for page_idx in range(len(content.pages)):
        text += (
            f'@@ Page {page_idx + 1} @@\n'
            + content.pages[page_idx].extract_text()
            + '\n\n'
        )
    print(text.strip())


@update_pwd_decorator
def parse_docx(file_path: str) -> None:
    """
    Parses the content of a DOCX file and prints it.

    Args:
        file_path: str: The path to the file to open.
    """
    print(f'[Reading DOCX file from {file_path}]')
    content = docx.Document(file_path)
    text = ''
    for i, para in enumerate(content.paragraphs):
        text += f'@@ Page {i + 1} @@\n' + para.text + '\n\n'
    print(text)


@update_pwd_decorator
def parse_latex(file_path: str) -> None:
    """
    Parses the content of a LaTex file and prints it.

    Args:
        file_path: str: The path to the file to open.
    """
    print(f'[Reading LaTex file from {file_path}]')
    with open(file_path) as f:
        data = f.read()
    text = LatexNodes2Text().latex_to_text(data)
    print(text.strip())


def _base64_img(file_path: str) -> str:
    with open(file_path, 'rb') as image_file:
        encoded_image = base64.b64encode(image_file.read()).decode('utf-8')
    return encoded_image


def _base64_video(file_path: str, frame_interval: int = 10) -> list[str]:
    import cv2

    video = cv2.VideoCapture(file_path)
    base64_frames = []
    frame_count = 0
    while video.isOpened():
        success, frame = video.read()
        if not success:
            break
        if frame_count % frame_interval == 0:
            _, buffer = cv2.imencode('.jpg', frame)
            base64_frames.append(base64.b64encode(buffer).decode('utf-8'))
        frame_count += 1
    video.release()
    return base64_frames


def _prepare_image_messages(task: str, base64_image: str):
    return [
        {
            'role': 'user',
            'content': [
                {'type': 'text', 'text': task},
                {
                    'type': 'image_url',
                    'image_url': {'url': f'data:image/jpeg;base64,{base64_image}'},
                },
            ],
        }
    ]


@update_pwd_decorator
def parse_audio(file_path: str, model: str = 'whisper-1') -> None:
    """
    Parses the content of an audio file and prints it.

    Args:
        file_path: str: The path to the audio file to transcribe.
        model: Optional[str]: The audio model to use for transcription. Defaults to 'whisper-1'.
    """
    print(f'[Transcribing audio file from {file_path}]')
    try:
        # TODO: record the COST of the API call
        with open(file_path, 'rb') as audio_file:
            transcript = client.audio.translations.create(model=model, file=audio_file)
        print(transcript.text)

    except Exception as e:
        print(f'Error transcribing audio file: {e}')


@update_pwd_decorator
def parse_image(
    file_path: str, task: str = 'Describe this image as detail as possible.'
) -> None:
    """
    Parses the content of an image file and prints the description.

    Args:
        file_path: str: The path to the file to open.
        task: Optional[str]: The task description for the API call. Defaults to 'Describe this image as detail as possible.'.
    """
    print(f'[Reading image file from {file_path}]')
    # TODO: record the COST of the API call
    try:
        base64_image = _base64_img(file_path)
        response = client.chat.completions.create(
            model=OPENAI_MODEL,
            messages=_prepare_image_messages(task, base64_image),
            max_tokens=MAX_TOKEN,
        )
        content = response.choices[0].message.content
        print(content)

    except Exception as error:
        print(f'Error with the request: {error}')


@update_pwd_decorator
def parse_video(
    file_path: str,
    task: str = 'Describe this image as detail as possible.',
    frame_interval: int = 30,
) -> None:
    """
    Parses the content of an image file and prints the description.

    Args:
        file_path: str: The path to the video file to open.
        task: Optional[str]: The task description for the API call. Defaults to 'Describe this image as detail as possible.'.
        frame_interval: Optional[int]: The interval between frames to analyze. Defaults to 30.

    """
    print(
        f'[Processing video file from {file_path} with frame interval {frame_interval}]'
    )

    task = task or 'This is one frame from a video, please summarize this frame.'
    base64_frames = _base64_video(file_path)
    selected_frames = base64_frames[::frame_interval]

    if len(selected_frames) > 30:
        new_interval = len(base64_frames) // 30
        selected_frames = base64_frames[::new_interval]

    print(f'Totally {len(selected_frames)} would be analyze...\n')

    idx = 0
    for base64_frame in selected_frames:
        idx += 1
        print(f'Process the {file_path}, current No. {idx * frame_interval} frame...')
        # TODO: record the COST of the API call
        try:
            response = client.chat.completions.create(
                model=OPENAI_MODEL,
                messages=_prepare_image_messages(task, base64_frame),
                max_tokens=MAX_TOKEN,
            )

            content = response.choices[0].message.content
            current_frame_content = f"Frame {idx}'s content: {content}\n"
            print(current_frame_content)

        except Exception as error:
            print(f'Error with the request: {error}')


@update_pwd_decorator
def parse_pptx(file_path: str) -> None:
    """
    Parses the content of a pptx file and prints it.

    Args:
        file_path: str: The path to the file to open.
    """
    print(f'[Reading PowerPoint file from {file_path}]')
    try:
        pres = Presentation(str(file_path))
        text = []
        for slide_idx, slide in enumerate(pres.slides):
            text.append(f'@@ Slide {slide_idx + 1} @@')
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text.append(shape.text)
        print('\n'.join(text))

    except Exception as e:
        print(f'Error reading PowerPoint file: {e}')


# def apply_git_patch(patch_str, directory='/workspace'):
#     """
#     Apply a git patch from a string to a designated folder within the given directory.

#     Parameters:
#     - patch_str (str): The git patch as a string.
#     - directory (str): The directory containing the subdirectory where the patch should be applied.

#     Returns:
#     - bool: True if the patch was successfully applied, False otherwise.
#     """
#     try:
#         # Ensure the target directory exists
#         if not os.path.isdir(directory):
#             raise FileNotFoundError(f'The directory {directory} does not exist.')

#         # Identify the full path to the subdirectory (assuming that there is only one subdirectory)
#         subdirectories = [
#             d
#             for d in os.listdir(directory)
#             if os.path.isdir(os.path.join(directory, d))
#         ]

#         if len(subdirectories) != 1:
#             raise FileNotFoundError(
#                 f'Expected exactly one subdirectory in {directory}, found {len(subdirectories)}.'
#             )

#         subdirectory_path = os.path.join(directory, subdirectories[0])
#         print(subdirectory_path)
#         # Change the current working directory to the subdirectory
#         os.chdir(subdirectory_path)

#         # Apply the patch
#         result = subprocess.run(
#             ['git', 'apply', '-'], input=patch_str, text=True, capture_output=True
#         )

#         # Check if the patch was applied successfully
#         if result.returncode != 0:
#             print(f'Error applying patch: {result.stderr}')
#             return False

#         print('Patch applied successfully.')
#         return True

#     except Exception as e:
#         print(f'An error occurred: {e}')
#         return False


# def try_apply_git_patch(patch_str, directory='./workspace'):
#     """Apply a git patch from a string to a designated folder, then revert it.

#     Parameters:
#     - patch_str (str): The git patch as a string.
#     - directory (str): The directory where the patch should be applied.

#     Returns:
#     - str: The text of the new file if the patch was successfully applied,
#            an error message otherwise.
#     """
#     import os
#     import subprocess

#     try:
#         # Ensure the target directory exists
#         if not os.path.isdir(directory):
#             raise FileNotFoundError(f'The directory {directory} does not exist.')

#         # Change the current working directory to the target directory
#         os.chdir(directory)

#         # Create a branch for safely applying the patch
#         subprocess.run(
#             ['git', 'checkout', '-b', 'temp_patch_branch'], capture_output=True
#         )

#         # Apply the patch
#         result = subprocess.run(
#             ['git', 'apply', '-'], input=patch_str, text=True, capture_output=True
#         )

#         # Check if the patch was applied successfully
#         if result.returncode != 0:
#             return f'Error applying patch: {result.stderr}'

#         # Get the status of the changes
#         status_result = subprocess.run(
#             ['git', 'status', '-s'], text=True, capture_output=True
#         )
#         if status_result.returncode != 0:
#             return f'Error getting git status: {status_result.stderr}'

#         # Get the diff of the applied patch
#         diff_result = subprocess.run(
#             ['git', 'diff', '--cached'], text=True, capture_output=True
#         )
#         if diff_result.returncode != 0:
#             return f'Error getting git diff: {diff_result.stderr}'

#         # Revert the applied patch
#         subprocess.run(['git', 'checkout', '-'], capture_output=True)
#         subprocess.run(
#             ['git', 'branch', '-D', 'temp_patch_branch'], capture_output=True
#         )

#         print('Patch applied and reverted successfully.')
#         return diff_result.stdout

#     except Exception as e:
#         return f'An error occurred: {e}'

# MAX_CONTEXT_LENGTH = 128000
# TOP_N = 3
# CONTEXT_WINDOW = 10
# ADD_SPACE = False
# STICKY_SCROLL = False
# NO_LINE_NUMBER = False
# TEMPERATURE = 0.8
# NUM_SAMPLES = 4
# LOC_INTERVAL = True
# FINE_GRAIN_LOC_ONLY = False
# COT = True
# DIFF_FORMAT = True
# STOP_AT_N_UNIQUE_VALID_SAMPLES = -1
# MAX_SAMPLES = 2
# SKIP_GREEDY = False

# def _parse_model_return_lines(content: str) -> list[str]:
#     if content:
#         return content.strip().split('\n')
#     return ['']

# def install_agentless_libraries():
#     import sys

#     def install_libcst():
#         print('libcst not found. Installing...')
#         subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'libcst'])
#         print('libcst installed successfully.')

#     install_libcst()

# def retrieve_structure(directory_path="./workspace"):
#     """Create the structure of the repository directory by parsing Python files.

#     If a pickled structure exists, load it. Otherwise, create it and pickle the result.

#     :param directory_path: Path to the repository directory.
#     :return: A dictionary representing the structure.
#     """

#     from Agentless.agentless.util.preprocess_data import (
#         filter_none_python,
#         filter_out_test_files,
#     )
#     from Agentless.get_repo_structure.get_repo_structure import (
#         parse_python_file,
#     )

#     install_agentless_libraries()

#     pickle_file = os.path.join(directory_path, 'repo_structure.pkl')

#     # Try to load the pickle file if it exists
#     if os.path.isfile(pickle_file):
#         with open(pickle_file, 'rb') as pf:
#             return pickle.load(pf)

#     # Proceed to create the structure if pickle file doesn't exist
#     structure: Dict[str, Any] = {}

#     for root, _, files in os.walk(directory_path):
#         repo_name = os.path.basename(directory_path)
#         relative_root = os.path.relpath(root, directory_path)

#         if relative_root == '.':
#             relative_root = repo_name

#         curr_struct = structure
#         for part in relative_root.split(os.sep):
#             if part not in curr_struct:
#                 curr_struct[part] = {}
#             curr_struct = curr_struct[part]

#         for file_name in files:
#             if file_name.endswith('.py'):
#                 file_path = os.path.join(root, file_name)
#                 class_info, function_names, file_lines = parse_python_file(
#                     file_path
#                 )
#                 curr_struct[file_name] = {
#                     'classes': class_info,
#                     'functions': function_names,
#                     'text': file_lines,
#                 }
#             else:
#                 curr_struct[file_name] = {}

#     filter_none_python(structure)
#     filter_out_test_files(structure)

#     # Save the structure to a pickle file
#     with open(pickle_file, 'wb') as pf:
#         pickle.dump(structure, pf)

#     return structure

# def agentless_file_localization(problem_statement):
#     from Agentless.agentless.fl.FL import LLMFL
#     from Agentless.agentless.util.preprocess_data import (
#         correct_file_paths,
#         get_full_file_paths_and_classes_and_functions,
#         show_project_structure,
#     )

#     install_agentless_libraries()

#     structure = retrieve_structure()

#     message = LLMFL.obtain_relevant_files_prompt.format(
#         problem_statement=problem_statement,
#         structure=show_project_structure(structure).strip(),
#     ).strip()
#     logger.info(f'prompting with message:\n{message}')
#     logger.info('=' * 80)

#     raw_output = client.chat.completions.create(
#         model=OPENAI_MODEL,
#         messages=[{'content': message, 'role': 'user'}],
#         temperature=0,
#         max_tokens=MAX_TOKEN
#     )['choices'][0]['message']['content']
#     model_found_files = _parse_model_return_lines(raw_output)

#     files, classes, functions = get_full_file_paths_and_classes_and_functions(
#         structure
#     )

#     # Never match partial paths for consistency with main Agentless results
#     found_files = correct_file_paths(model_found_files, files, False)
#     logger.info(found_files)

#     return "\n".join(found_files)

# def agentless_related_localization(problem_statement, found_files):
#     from Agentless.agentless.fl.FL import LLMFL
#     from Agentless.agentless.util.api_requests import (
#         num_tokens_from_messages,
#     )
#     from Agentless.agentless.util.compress_file import get_skeleton
#     from Agentless.agentless.util.postprocess_data import (
#         extract_code_blocks,
#         extract_locs_for_files,
#     )
#     from Agentless.agentless.util.preprocess_data import (
#         get_repo_files,
#     )

#     structure = retrieve_structure()

#     if len(found_files) != 0:
#         file_names = found_files[:TOP_N]
#     else:
#         file_names = found_files
#     file_contents = get_repo_files(structure, file_names)
#     compressed_file_contents = {
#         fn: get_skeleton(code) for fn, code in file_contents.items()
#     }
#     contents = [
#         LLMFL.file_content_in_block_template.format(file_name=fn, file_content=code)
#         for fn, code in compressed_file_contents.items()
#     ]
#     file_contents = ''.join(contents)
#     template = (
#         LLMFL.obtain_relevant_functions_and_vars_from_compressed_files_prompt_more
#     )
#     message = template.format(
#         problem_statement=problem_statement, file_contents=file_contents
#     )

#     def message_too_long(message):
#         return num_tokens_from_messages(message, OPENAI_MODEL) >= MAX_CONTEXT_LENGTH

#     while message_too_long(message) and len(contents) > 1:
#         logger.info(f'reducing to \n{len(contents)} files')
#         contents = contents[:-1]
#         file_contents = ''.join(contents)
#         message = template.format(
#             problem_statement=problem_statement, file_contents=file_contents
#         )  # Recreate message

#     if message_too_long(message):
#         logger.info
#         raise ValueError(
#             'The remaining file content is too long to fit within the context length'
#         )
#     logger.info(f'prompting with message:\n{message}')
#     logger.info('=' * 80)

#     raw_output = client.chat.completions.create(
#         model=OPENAI_MODEL,
#         messages=[{'content': message, 'role': 'user'}],
#         temperature=0,
#         max_tokens=MAX_TOKEN,
#     )['choices'][0]['message']['content']

#     model_found_locs = extract_code_blocks(raw_output)
#     model_found_locs_separated = extract_locs_for_files(
#         model_found_locs, file_names
#     )

#     logger.info('==== raw output ====')
#     logger.info(raw_output)
#     logger.info('=' * 80)
#     logger.info('==== extracted locs ====')
#     for loc in model_found_locs_separated:
#         logger.info(loc)
#     logger.info('=' * 80)

#     print(raw_output)

#     return model_found_locs_separated

# def agentless_line_level_localization(
#     file_names, found_related_locs, problem_statement
# ):

#     structure = retrieve_structure()

#     from Agentless.agentless.fl.FL import LLMFL
#     from Agentless.agentless.repair.repair import (
#         construct_topn_file_context,
#     )
#     from Agentless.agentless.util.api_requests import (
#         num_tokens_from_messages,
#     )
#     from Agentless.agentless.util.postprocess_data import (
#         extract_code_blocks,
#         extract_locs_for_files,
#     )
#     from Agentless.agentless.util.preprocess_data import (
#         get_repo_files,
#     )

#     coarse_locs = {}
#     for i, pred_file in enumerate(file_names):
#         if len(found_related_locs) > i:
#             coarse_locs[pred_file] = found_related_locs[i]

#     file_contents = get_repo_files(structure, file_names)
#     topn_content, _ = construct_topn_file_context(
#         coarse_locs,
#         file_names,
#         file_contents,
#         structure,
#         context_window=CONTEXT_WINDOW,
#         loc_interval=True,
#         add_space=ADD_SPACE,
#         sticky_scroll=STICKY_SCROLL,
#         no_line_number=NO_LINE_NUMBER,
#     )
#     if NO_LINE_NUMBER:
#         template = LLMFL.obtain_relevant_code_combine_top_n_no_line_number_prompt
#     else:
#         template = LLMFL.obtain_relevant_code_combine_top_n_prompt
#     message = template.format(
#         problem_statement=problem_statement, file_contents=topn_content
#     )
#     logger.info(f'prompting with message:\n{message}')
#     logger.info('=' * 80)
#     assert num_tokens_from_messages(message, OPENAI_MODEL) < MAX_CONTEXT_LENGTH

#     raw_trajs = client.chat.completions.create(
#         model=OPENAI_MODEL,
#         messages=[{'content': message, 'role': 'user'}],
#         temperature=TEMPERATURE,
#         max_tokens=OPENAI_MODEL,
#         n=NUM_SAMPLES,
#     )['choices']
#     # raw_trajs = model.codegen(message, num_samples=NUM_SAMPLES)

#     # Merge trajectories
#     raw_outputs = [raw_traj['message']['content'] for raw_traj in raw_trajs]

#     model_found_locs_separated_in_samples = []
#     for raw_output in raw_outputs:
#         model_found_locs = extract_code_blocks(raw_output)
#         model_found_locs_separated = extract_locs_for_files(
#             model_found_locs, file_names
#         )
#         model_found_locs_separated_in_samples.append(model_found_locs_separated)

#         logger.info('==== raw output ====')
#         logger.info(raw_output)
#         logger.info('=' * 80)
#         print(raw_output)
#         print('=' * 80)
#         logger.info('==== extracted locs ====')
#         for loc in model_found_locs_separated:
#             logger.info(loc)
#         logger.info('=' * 80)
#     logger.info('Localizations: ')
#     logger.info(model_found_locs_separated_in_samples)
#     logger.info('==== Input coarse_locs')
#     coarse_info = ''
#     for fn, found_locs in coarse_locs.items():
#         coarse_info += f'### {fn}\n'
#         if isinstance(found_locs, str):
#             coarse_info += found_locs + '\n'
#         else:
#             coarse_info += '\n'.join(found_locs) + '\n'
#     logger.info('\n' + coarse_info)
#     if len(model_found_locs_separated_in_samples) == 1:
#         model_found_locs_separated_in_samples = (
#             model_found_locs_separated_in_samples[0]
#         )

#     merged_locs = []
#     first_merge = (
#         model_found_locs_separated_in_samples[0]
#         + model_found_locs_separated_in_samples[1]
#     )
#     second_merge = (
#         model_found_locs_separated_in_samples[2]
#         + model_found_locs_separated_in_samples[3]
#     )
#     merged_locs.append(first_merge)
#     merged_locs.append(second_merge)

#     def flatten_innermost_strings(nested_list):
#         if isinstance(nested_list, list):
#             # Check if we have a list with a single item which is a string
#             if len(nested_list) == 1 and isinstance(nested_list[0], str):
#                 return nested_list[0]
#             # Otherwise, recursively process each element in the list
#             return [flatten_innermost_strings(item) for item in nested_list]
#         else:
#             return nested_list

#     return flatten_innermost_strings(merged_locs)

# def agentless_repair(
#     found_files, both_found_edit_locs, problem_statement
# ):

#     structure = retrieve_structure()

#     from Agentless.agentless.repair.repair import (
#         _post_process_multifile_repair,
#         construct_topn_file_context,
#         repair_prompt_combine_topn,
#         repair_prompt_combine_topn_cot,
#         repair_prompt_combine_topn_cot_diff,
#         repair_relevant_file_instruction,
#     )
#     from Agentless.agentless.util.preprocess_data import (
#         get_full_file_paths_and_classes_and_functions,
#     )

#     for found_edit_locs in both_found_edit_locs:
#         pred_files = found_files[:TOP_N]
#         file_to_edit_locs = {}
#         for i, pred_file in enumerate(pred_files):
#             if len(found_edit_locs) > i:
#                 file_to_edit_locs[pred_file] = found_edit_locs[i]

#         files, _, _ = get_full_file_paths_and_classes_and_functions(structure)

#         # Construct file contents
#         file_contents = dict()
#         for i, pred_file in enumerate(pred_files):
#             content = None

#             for file_content in files:
#                 if file_content[0] == pred_file:
#                     content = '\n'.join(file_content[1])
#                     file_contents[pred_file] = content
#                     break

#         logger.info('FILE CONTENTS')
#         logger.info(file_contents)

#         topn_content, file_loc_intervals = construct_topn_file_context(
#             file_to_edit_locs,
#             pred_files,
#             file_contents,
#             structure,
#             context_window=CONTEXT_WINDOW,
#             loc_interval=LOC_INTERVAL,
#             fine_grain_loc_only=FINE_GRAIN_LOC_ONLY,
#             add_space=ADD_SPACE,
#             no_line_number=NO_LINE_NUMBER,
#             sticky_scroll=STICKY_SCROLL,
#         )

#         if topn_content.strip() == '':
#             return []

#         prompt_template = (
#             repair_prompt_combine_topn_cot_diff
#             if COT and DIFF_FORMAT
#             else repair_prompt_combine_topn_cot
#             if COT
#             else repair_prompt_combine_topn
#         )
#         file_instruction = repair_relevant_file_instruction
#         message = prompt_template.format(
#             repair_relevant_file_instruction=file_instruction,
#             problem_statement=problem_statement,
#             content=topn_content.rstrip(),
#         ).strip()
#         logger.info(f'prompting with message:\n{message}')

#         all_generations, counts, _, prev_contents, file_names = [], [], [], [], []
#         sample_responses: List[Dict[str, Any]] = []
#         # Using early stopping will cost more since the input tokens will be charged multiple times.
#         # For now we disable it.
#         assert STOP_AT_N_UNIQUE_VALID_SAMPLES == -1

#         if SKIP_GREEDY:
#             greedy_traj = {
#                 'response': '',
#                 'usage': {
#                     'completion_tokens': 0,
#                     'prompt_tokens': 0,
#                 },
#             }
#         else:
#             greedy_traj = client.chat.completions.create(
#                 model=OPENAI_MODEL,
#                 messages=[{'content': message, 'role': 'user'}],
#                 temperature=0,
#                 max_tokens=1024,
#             )['choices'][0]
#         sample_responses.append(greedy_traj)

#         if MAX_SAMPLES - 1:
#             sample_trajs = client.chat.completions.create(
#                 model=OPENAI_MODEL,
#                 messages=[{'content': message, 'role': 'user'}],
#                 temperature=TEMPERATURE,
#                 max_tokens=1024,
#                 n=MAX_SAMPLES - 1,
#             )['choices']
#         else:
#             sample_trajs = []

#         sample_responses.extend(sample_trajs)

#     count = 0
#     raw_outputs = []
#     while count < MAX_SAMPLES:
#         print(f'trying the {count + 1}-th sample ...')
#         raw_output = sample_responses[count]['message']['content']
#         count += 1

#         all_generations.append(raw_output)

#         edited_file, new_content = _post_process_multifile_repair(
#             raw_output,
#             file_contents,
#             logger,
#             file_loc_intervals,
#             diff_format=DIFF_FORMAT,
#         )

#         if new_content == '':
#             prev_contents.append('')
#             file_names.append('')
#         else:
#             prev_content = file_contents[edited_file]
#             prev_contents.append(prev_content)
#             file_names.append(edited_file)

#         counts.append(count)
#         raw_outputs.append(raw_output)
#     return raw_outputs, prev_contents, file_names

# def post_process_raw_output(
#     raw_output_text, file_contents, file_loc_intervals
# ):
#     "post_process_raw_output"
#     from difflib import unified_diff

#     from Agentless.agentless.repair.repair import (
#         _post_process_multifile_repair,
#     )
#     from Agentless.agentless.util.postprocess_data import (
#         check_code_differ_by_just_empty_lines,
#         check_syntax,
#         fake_git_repo,
#         lint_code,
#     )

#     git_diffs = ''
#     raw_git_diffs = ''
#     lint_success = False
#     content = ''
#     try:
#         edited_file, new_content = _post_process_multifile_repair(
#             raw_output_text,
#             file_contents,
#             logger,
#             file_loc_intervals,
#             diff_format=DIFF_FORMAT,
#         )
#         logger.info('EDITED FILE')
#         logger.info(edited_file)
#         logger.info('NEW CONTENT')
#         logger.info(new_content)
#         if (edited_file in file_contents) and new_content:
#             content = file_contents[edited_file]

#             git_diff = fake_git_repo(
#                 'playground', edited_file, content, new_content
#             )
#             logger.info('1 git diff')
#             logger.info(git_diff)

#             raw_git_diffs += '\n' + git_diff.replace(
#                 '\\ No newline at end of file\n', ''
#             )

#             syntax_success = check_syntax(new_content)
#             lint_success, prev_errors, errors = lint_code(
#                 'playground', 'test.py', new_content, file_contents[edited_file]
#             )

#             differ_by_empty_lines = check_code_differ_by_just_empty_lines(
#                 new_content, file_contents[edited_file]
#             )

#             print(lint_success, prev_errors, errors, differ_by_empty_lines)
#             logger.info(
#                 f'checks pass: {lint_success, prev_errors, errors, differ_by_empty_lines}'
#             )
#             if syntax_success and not differ_by_empty_lines:
#                 git_diffs = raw_git_diffs
#             else:
#                 git_diffs = ''  # no need to evaluate
#         else:
#             diff = list(
#                 unified_diff(
#                     content.split('\n'),
#                     new_content.split('\n'),
#                     fromfile=edited_file,
#                     tofile=edited_file,
#                     lineterm='',
#                 )
#             )
#             print('Failed parsing diff!')
#             print('\n'.join(diff))
#     except Exception as e:
#         print(raw_output_text)
#         print(e)

#     return git_diffs, raw_git_diffs, content

# def agentless_post_process_repair(
#     pred_files,
#     both_found_edit_locs,
#     generation_original_file_contents,
#     generation_pred_files,
#     raw_outputs,
# ):
#     """
#     apply some diff formatting. Also strips comments and trailing spaces to better identify identical patches
#     """
#     from Agentless.agentless.util.postprocess_data import normalize_patch
#     from Agentless.agentless.util.preprocess_data import (
#         transfer_arb_locs_to_locs,
#     )

#     num_each_edit_loc = int(len(raw_outputs) / 2)

#     processed_patches = []
#     for generation_idx, raw_output_text in enumerate(raw_outputs):
#         if num_each_edit_loc > generation_idx:
#             found_edit_locs = both_found_edit_locs[0]
#         else:
#             found_edit_locs = both_found_edit_locs[1]
#         if raw_output_text != '':
#             try:
#                 original_file_content = generation_original_file_contents[
#                     generation_idx
#                 ]
#                 pred_file = generation_pred_files[
#                     generation_idx
#                 ]  # Not sure if this works

#                 git_diffs = ''

#                 file_contents = {pred_file: original_file_content}

#                 file_loc_intervals = dict()

#                 for i, tmp_pred_file in enumerate(pred_files):
#                     if tmp_pred_file != pred_file:
#                         continue
#                     if len(found_edit_locs) > i:
#                         # try:
#                         _, context_intervals = transfer_arb_locs_to_locs(
#                             found_edit_locs[i],
#                             None,
#                             pred_files[i],
#                             CONTEXT_WINDOW,
#                             LOC_INTERVAL,
#                             FINE_GRAIN_LOC_ONLY,
#                             file_content=file_contents[pred_file]
#                             if pred_file in file_contents
#                             else '',
#                         )
#                         logger.info('context interval')
#                         logger.info(context_intervals)
#                     else:
#                         _, context_intervals = [], []  # default values.

#             except Exception as e:
#                 logger.info('file loc interval error')
#                 logger.info(e)
#                 print(e)
#                 raw_output_text = ''

#             file_loc_intervals[pred_file] = context_intervals

#             logger.info('RAW OUTPUT TEXT')
#             logger.info(raw_output_text)

#             if raw_output_text:
#                 logger.info('FILE LOCAL INTERVAL')
#                 logger.info(file_loc_intervals)
#                 logger.info('FILE CONTENT')
#                 logger.info(file_contents)
#                 git_diffs, raw_git_diffs, content = post_process_raw_output(
#                     raw_output_text, file_contents, logger, file_loc_intervals
#                 )

#                 # Setting 0 as the instance_id since this function doesn't
#                 # actually use the instance_id for anything
#                 logger.info('GIT DIFF BEFORE NORMALIZING')
#                 logger.info(git_diffs)
#                 patch_lines = git_diffs.split('\n')
#                 patch_lines = [
#                     line for line in patch_lines if not line.startswith('index')
#                 ]
#                 git_diffs = '\n'.join(patch_lines)

#                 normalized_patch = normalize_patch(
#                     0, git_diffs, original_file_content
#                 )

#                 normalized_lines = normalized_patch.split('\n')
#                 normalized_lines = [
#                     line
#                     for line in normalized_lines
#                     if not line.startswith('index')
#                 ]
#                 normalized_patch = '\n'.join(patch_lines)

#                 if normalized_patch.lstrip():
#                     processed_patches.append(
#                         [
#                             normalized_patch,
#                             git_diffs.replace(
#                                 '\\ No newline at end of file\n', ''
#                             ).lstrip(),
#                         ]
#                     )

#     return [patch for patch in processed_patches if patch[0].strip() != '']


# def most_frequent_string(processed_patches):
#     "fdkfjdkf"
#     if not processed_patches:
#         return None  # Return None if the list is empty

#     # Count the frequency of each string in the list
#     counter = Counter([patch[0] for patch in processed_patches])

#     # Find the string with the highest frequency
#     most_common_normalized_patch, freq = counter.most_common(1)[0]
#     print('Number of times most frequent patch occured: ', freq)
#     for patch in processed_patches:
#         if patch[0] == most_common_normalized_patch:
#             return patch[1]

# def test_output(test):
#     """test how observations work
#         Apply a git patch from a string to a designated folder within the given directory.

#     Parameters:
#     - patch_str (str): The git patch as a string.
#     - directory (str): The directory containing the subdirectory where the patch should be applied.

#     Returns:
#     - bool: True if the patch was successfully applied, False otherwise.
#     """
#     # print("Hello there")
#     return("Hello there!")

__all__ = [
    # file operation
    'open_file',
    'goto_line',
    'scroll_down',
    'scroll_up',
    'create_file',
    'edit_file',
    'insert_content_at_line',
    'search_dir',
    'search_file',
    'find_file',
    # readers
    'parse_pdf',
    'parse_docx',
    'parse_latex',
    'parse_pptx',
    # Agentless skills
    # 'test_output',
    # 'apply_git_patch',
]

if OPENAI_API_KEY and OPENAI_BASE_URL:
    __all__ += ['parse_audio', 'parse_video', 'parse_image']

DOCUMENTATION = ''
for func_name in __all__:
    func = globals()[func_name]

    cur_doc = func.__doc__
    # remove indentation from docstring and extra empty lines
    cur_doc = '\n'.join(filter(None, map(lambda x: x.strip(), cur_doc.split('\n'))))
    # now add a consistent 4 indentation
    cur_doc = '\n'.join(map(lambda x: ' ' * 4 + x, cur_doc.split('\n')))

    fn_signature = f'{func.__name__}' + str(signature(func))
    DOCUMENTATION += f'{fn_signature}:\n{cur_doc}\n\n'
